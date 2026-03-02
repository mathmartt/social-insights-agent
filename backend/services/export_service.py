"""
ExportService: generates PowerPoint (.pptx) and CSV exports.
"""
import io
import csv
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ─── Google colour palette ────────────────────────────────────────────────────
BLUE   = RGBColor(0x42, 0x85, 0xF4)
GREEN  = RGBColor(0x34, 0xA8, 0x53)
YELLOW = RGBColor(0xFB, 0xBC, 0x04)
RED    = RGBColor(0xEA, 0x43, 0x35)
DARK   = RGBColor(0x20, 0x21, 0x24)
GREY   = RGBColor(0x5F, 0x63, 0x68)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF8, 0xF9, 0xFA)


def _add_slide(prs: Presentation, layout_idx: int = 6):
    """Blank layout slide."""
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _txbox(slide, text: str, left, top, width, height,
           font_size=14, bold=False, colour=DARK, align=PP_ALIGN.LEFT,
           wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return txBox


def _rect(slide, left, top, width, height, fill_colour):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    return shape


def _footer(slide, text: str):
    """Add footer line at bottom of slide."""
    _txbox(slide, text,
           Inches(0.4), Inches(7.2), Inches(9.2), Inches(0.4),
           font_size=8, colour=GREY, align=PP_ALIGN.CENTER)


# ─── Slide builders ───────────────────────────────────────────────────────────

def _slide_title(prs, query: str, date_range: dict, total: int, relevant: int):
    slide = _add_slide(prs)
    # Background accent bar
    _rect(slide, Inches(0), Inches(0), Inches(10), Inches(1.2), BLUE)

    # Header label
    _txbox(slide, "Social Insights Agent  ·  Google Marketing",
           Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
           font_size=13, bold=False, colour=WHITE)

    # Big query
    _txbox(slide, f'"{query}"',
           Inches(0.5), Inches(1.5), Inches(9), Inches(2.5),
           font_size=26, bold=True, colour=DARK, wrap=True)

    # Meta row
    date_str = f"{date_range.get('start','?')}  →  {date_range.get('end','?')}"
    _txbox(slide, date_str,
           Inches(0.5), Inches(4.1), Inches(4), Inches(0.5),
           font_size=13, colour=GREY)
    _txbox(slide, f"{total:,} comments analysed  ·  {relevant:,} relevant",
           Inches(0.5), Inches(4.6), Inches(5), Inches(0.5),
           font_size=13, colour=GREY)

    _footer(slide, "Source: Google Owned Social Handles via Sprinklr")
    return slide


def _slide_sentiment(prs, sentiment: dict, trend: list):
    slide = _add_slide(prs)
    _rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.7), BLUE)
    _txbox(slide, "Sentiment Overview",
           Inches(0.4), Inches(0.1), Inches(9), Inches(0.5),
           font_size=16, bold=True, colour=WHITE)

    pos = sentiment.get("positive_pct", 0)
    neu = sentiment.get("neutral_pct", 0)
    neg = sentiment.get("negative_pct", 0)
    total_rel = sentiment.get("total_relevant", 0)

    # Three big numbers
    for i, (label, pct, colour) in enumerate([
        ("Positive", pos, GREEN),
        ("Neutral",  neu, BLUE),
        ("Negative", neg, RED),
    ]):
        x = Inches(0.5 + i * 3.2)
        _rect(slide, x, Inches(1.0), Inches(2.8), Inches(1.8), LIGHT)
        _txbox(slide, f"{pct:.0f}%",
               x + Inches(0.1), Inches(1.1), Inches(2.6), Inches(1.0),
               font_size=44, bold=True, colour=colour, align=PP_ALIGN.CENTER)
        _txbox(slide, label,
               x + Inches(0.1), Inches(2.05), Inches(2.6), Inches(0.4),
               font_size=12, colour=GREY, align=PP_ALIGN.CENTER)

    _txbox(slide, f"Based on {total_rel:,} relevant comments",
           Inches(0.5), Inches(3.0), Inches(9), Inches(0.4),
           font_size=11, colour=GREY)

    # Trend table (text)
    if trend:
        _txbox(slide, "Sentiment Trend (by month)",
               Inches(0.5), Inches(3.5), Inches(9), Inches(0.4),
               font_size=12, bold=True, colour=DARK)
        rows_text = "  ".join(
            f"{t['period']}: +{t['positive']} /{t['neutral']} -{t['negative']}"
            for t in trend[:8]
        )
        _txbox(slide, rows_text,
               Inches(0.5), Inches(3.9), Inches(9), Inches(2.0),
               font_size=10, colour=GREY, wrap=True)

    _footer(slide, "Source: Google Owned Social Handles via Sprinklr")


def _slide_findings(prs, summary: str, themes: list):
    slide = _add_slide(prs)
    _rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.7), GREEN)
    _txbox(slide, "Key Findings",
           Inches(0.4), Inches(0.1), Inches(9), Inches(0.5),
           font_size=16, bold=True, colour=WHITE)

    _txbox(slide, summary,
           Inches(0.5), Inches(1.0), Inches(9), Inches(3.0),
           font_size=13, colour=DARK, wrap=True)

    if themes:
        _txbox(slide, "Key Themes",
               Inches(0.5), Inches(4.2), Inches(9), Inches(0.4),
               font_size=12, bold=True, colour=DARK)
        themes_text = "  ·  ".join(themes)
        _txbox(slide, themes_text,
               Inches(0.5), Inches(4.7), Inches(9), Inches(1.0),
               font_size=11, colour=GREY, wrap=True)

    _footer(slide, "Source: Google Owned Social Handles via Sprinklr")


def _slide_actions(prs, action_items: list):
    slide = _add_slide(prs)
    _rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.7), YELLOW)
    _txbox(slide, "Recommended Actions",
           Inches(0.4), Inches(0.1), Inches(9), Inches(0.5),
           font_size=16, bold=True, colour=DARK)

    colours = [BLUE, GREEN, RED, YELLOW]
    for i, item in enumerate(action_items[:4]):
        y = Inches(0.9 + i * 1.55)
        c = colours[i % len(colours)]
        _rect(slide, Inches(0.4), y, Inches(0.07), Inches(1.2), c)
        priority = item.get("priority", "").upper()
        platform = item.get("platform", "")
        meta = f"{platform}  ·  {priority}" if platform else priority
        _txbox(slide, item.get("title", ""),
               Inches(0.65), y, Inches(8.8), Inches(0.45),
               font_size=12, bold=True, colour=DARK)
        _txbox(slide, item.get("description", ""),
               Inches(0.65), y + Inches(0.42), Inches(8.3), Inches(0.65),
               font_size=10, colour=GREY, wrap=True)
        _txbox(slide, meta,
               Inches(0.65), y + Inches(1.05), Inches(8.3), Inches(0.3),
               font_size=9, colour=c)

    _footer(slide, "Source: Google Owned Social Handles via Sprinklr")


def _slide_evidence(prs, comments: list):
    slide = _add_slide(prs)
    _rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.7), RED)
    _txbox(slide, "Comment Evidence",
           Inches(0.4), Inches(0.1), Inches(9), Inches(0.5),
           font_size=16, bold=True, colour=WHITE)

    pos_comments = [c for c in comments if c.get("comment_sentiment") == "Positive"][:4]
    neu_comments = [c for c in comments if c.get("comment_sentiment") == "Neutral"][:3]
    neg_comments = [c for c in comments if c.get("comment_sentiment") == "Negative"][:4]

    col_configs = [
        ("Positive ✓", pos_comments, GREEN,  Inches(0.3)),
        ("Neutral",    neu_comments, BLUE,   Inches(3.55)),
        ("Negative ✗", neg_comments, RED,    Inches(6.8)),
    ]

    for header, items, colour, x in col_configs:
        _txbox(slide, header, x, Inches(0.85), Inches(2.9), Inches(0.35),
               font_size=11, bold=True, colour=colour)
        for j, c in enumerate(items[:4]):
            y = Inches(1.25 + j * 1.55)
            _rect(slide, x, y, Inches(2.9), Inches(1.4), LIGHT)
            text = c.get("comment_text", "")[:150]
            if len(c.get("comment_text", "")) > 150:
                text += "…"
            _txbox(slide, text, x + Inches(0.08), y + Inches(0.05),
                   Inches(2.74), Inches(1.0),
                   font_size=8, colour=DARK, wrap=True)
            meta = f"{c.get('social_network','')} · {str(c.get('comment_date',''))[:10]}"
            _txbox(slide, meta, x + Inches(0.08), y + Inches(1.1),
                   Inches(2.74), Inches(0.28),
                   font_size=7, colour=GREY)

    _footer(slide, "Source: Google Owned Social Handles via Sprinklr")


# ─── Public API ──────────────────────────────────────────────────────────────

class ExportService:

    def create_pptx_bytes(self, data: dict) -> bytes:
        """
        Build a complete 5-slide PowerPoint and return it as bytes.
        """
        prs = Presentation()
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(7.5)

        query      = data.get("query", "Social Insights Report")
        time_range = data.get("time_range_covered", data.get("date_range", {}))
        total      = data.get("total_analyzed", 0)
        relevant   = data.get("sentiment_breakdown", {}).get("total_relevant", 0)
        sentiment  = data.get("sentiment_breakdown", {})
        trend      = data.get("sentiment_trend", [])
        summary    = data.get("written_summary", "")
        themes     = data.get("key_themes", [])
        actions    = data.get("action_items", [])
        comments   = data.get("relevant_comments", [])

        _slide_title(prs, query, time_range, total, relevant)
        _slide_sentiment(prs, sentiment, trend)
        _slide_findings(prs, summary, themes)
        _slide_actions(prs, actions)
        _slide_evidence(prs, comments)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def create_csv_bytes(self, comments: list[dict]) -> bytes:
        """
        Serialise the relevant comments list to CSV bytes.
        """
        if not comments:
            return b"No relevant comments found.\n"

        fields = [
            "comment_date", "social_network", "google_account",
            "campaign_name", "comment_text", "comment_sentiment",
            "comment_likes", "comment_author", "post_caption",
        ]

        buf = io.StringIO()
        writer = csv.DictWriter(buf, fieldnames=fields, extrasaction="ignore",
                                lineterminator="\n")
        writer.writeheader()
        writer.writerows(comments)
        return buf.getvalue().encode("utf-8")
